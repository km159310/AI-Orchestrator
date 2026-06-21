#!/usr/bin/env python3
"""Generates AI-Orchestrator-Features.pptx — a feature-overview deck for the
AI Orchestrator project. Run:  python make_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette (matches the orchestrator's light/professional theme) ──
BRAND_BLUE   = RGBColor(0x25, 0x63, 0xEB)
BRAND_INDIGO = RGBColor(0x4F, 0x46, 0xE5)
BRAND_PURPLE = RGBColor(0x7C, 0x3A, 0xED)
BRAND_GREEN  = RGBColor(0x05, 0x96, 0x69)
BRAND_AMBER  = RGBColor(0xD9, 0x77, 0x06)
BRAND_RED    = RGBColor(0xDC, 0x26, 0x26)
TEXT_DARK    = RGBColor(0x0F, 0x17, 0x2A)
TEXT_MID     = RGBColor(0x47, 0x55, 0x69)
TEXT_LIGHT   = RGBColor(0x94, 0xA3, 0xB8)
BG_LIGHT     = RGBColor(0xF6, 0xF8, 0xFB)
BG_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BORDER       = RGBColor(0xE5, 0xE9, 0xF0)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BLANK = prs.slide_layouts[6]


def _add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def _add_rounded(slide, x, y, w, h, fill, line=None, radius=0.04):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = radius
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def _add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT_DARK,
              align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    if "\n" in text:
        lines = text.split("\n")
        run = p.add_run(); run.text = lines[0]
        for line in lines[1:]:
            new_p = tf.add_paragraph()
            new_p.alignment = align
            r = new_p.add_run(); r.text = line
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    else:
        run = p.add_run(); run.text = text
    for para in tf.paragraphs:
        for r in para.runs:
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return tb


def _add_bullets(slide, x, y, w, h, bullets, *, size=14, color=TEXT_DARK, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run(); r.text = "• " + b
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = font
    return tb


def _add_footer(slide, idx, total):
    _add_text(slide, Inches(0.5), Inches(7.0), Inches(8), Inches(0.4),
              "AI Orchestrator — Feature Overview · ABC Bank Application",
              size=10, color=TEXT_LIGHT)
    _add_text(slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4),
              f"{idx} / {total}", size=10, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)


def _page_bg(slide, color=BG_LIGHT):
    bg = _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)
    return bg


def _header(slide, kicker, title, tagline=None, gradient_color=BRAND_BLUE):
    # Top accent bar
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.18), gradient_color)
    if kicker:
        _add_text(slide, Inches(0.5), Inches(0.45), Inches(8), Inches(0.35),
                  kicker.upper(), size=10, bold=True, color=gradient_color)
    _add_text(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.7),
              title, size=32, bold=True, color=TEXT_DARK)
    if tagline:
        _add_text(slide, Inches(0.5), Inches(1.50), Inches(12), Inches(0.4),
                  tagline, size=14, color=TEXT_MID)


def _card(slide, x, y, w, h, title, body_lines, *, accent=BRAND_BLUE):
    _add_rounded(slide, x, y, w, h, BG_WHITE, line=BORDER, radius=0.05)
    # accent strip on top
    _add_rect(slide, x + Emu(50000), y + Emu(50000), w - Emu(100000), Inches(0.06), accent)
    _add_text(slide, x + Inches(0.25), y + Inches(0.22), w - Inches(0.5), Inches(0.45),
              title, size=15, bold=True, color=TEXT_DARK)
    _add_bullets(slide, x + Inches(0.25), y + Inches(0.75), w - Inches(0.5), h - Inches(1.0),
                 body_lines, size=11.5, color=TEXT_MID)


def _phase_chip(slide, x, y, w, h, label, color):
    _add_rounded(slide, x, y, w, h, color, radius=0.4)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = Inches(0.18); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = BG_WHITE; r.font.name = "Calibri"


# ──────────────────────────────────────────────────────────────────
# SLIDES
# ──────────────────────────────────────────────────────────────────

SLIDES = []  # filled below for footer numbering


# ── Slide 1: title ────────────────────────────────────────────────
def slide_title(idx, total):
    s = prs.slides.add_slide(BLANK)
    # Full-bleed gradient (simulated via 3 layered rectangles)
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BRAND_BLUE)
    _add_rect(s, Inches(6), 0, SLIDE_W - Inches(6), SLIDE_H, BRAND_INDIGO)
    _add_rect(s, Inches(10), 0, SLIDE_W - Inches(10), SLIDE_H, BRAND_PURPLE)
    # Logo tile
    _add_rounded(s, Inches(1.0), Inches(2.6), Inches(1.0), Inches(1.0), BG_WHITE, radius=0.18)
    _add_text(s, Inches(1.0), Inches(2.78), Inches(1.0), Inches(0.7),
              "AI", size=36, bold=True, color=BRAND_BLUE, align=PP_ALIGN.CENTER)
    # Title
    _add_text(s, Inches(2.4), Inches(2.4), Inches(10), Inches(1.0),
              "AI Orchestrator", size=54, bold=True, color=BG_WHITE)
    _add_text(s, Inches(2.4), Inches(3.4), Inches(10), Inches(0.5),
              "Feature Overview", size=24, color=BG_WHITE)
    _add_text(s, Inches(2.4), Inches(4.1), Inches(10), Inches(0.5),
              "AI-Powered SDLC Pipeline · ABC Bank Application", size=14, color=BG_WHITE)
    _add_text(s, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
              "2026-06-06", size=10, color=BG_WHITE)


# ── Slide 2: What is it ───────────────────────────────────────────
def slide_what(idx, total):
    s = prs.slides.add_slide(BLANK)
    _page_bg(s)
    _header(s, "Slide 1 of " + str(total - 1), "What is AI Orchestrator?",
            "An end-to-end, AI-driven SDLC pipeline that takes a plain-text BRD and produces a running application.")
    _card(s, Inches(0.5),  Inches(2.3), Inches(4.0), Inches(4.3), "Input",
          ["Paste or upload a BRD (TXT / MD)",
           "Heuristic NLP extracts requirements",
           "Detects priority, ports, risks, stakeholders",
           "Stack picker: Java / Python / .NET"], accent=BRAND_BLUE)
    _card(s, Inches(4.7),  Inches(2.3), Inches(4.0), Inches(4.3), "Pipeline",
          ["7 sequential phases with AI agents",
           "Phase agents run one after another",
           "Human sign-off at each gate",
           "Reject with reason · verify · proceed",
           "Auto-advance on approval"], accent=BRAND_INDIGO)
    _card(s, Inches(8.9),  Inches(2.3), Inches(4.0), Inches(4.3), "Output",
          ["Real running bank app on disk",
           "Multi-port deployment (:3001, :3002)",
           "Login, dashboard, balance, transactions",
           "Jenkins pipeline trigger (or mock)",
           "Generated docs: SRS, API, ERD"], accent=BRAND_PURPLE)
    _add_footer(s, idx, total)


# ── Slide 3: 7-phase pipeline ─────────────────────────────────────
def slide_pipeline(idx, total):
    s = prs.slides.add_slide(BLANK)
    _page_bg(s)
    _header(s, "PIPELINE", "Seven AI-driven phases",
            "Each phase has its own agent crew, hero banner, and human sign-off.")
    phases = [
        ("Requirements", BRAND_BLUE),
        ("Design", BRAND_INDIGO),
        ("Development", RGBColor(0x08, 0x91, 0xB2)),
        ("Testing", BRAND_GREEN),
        ("Security", BRAND_RED),
        ("Deployment", BRAND_PURPLE),
        ("Review", BRAND_AMBER),
    ]
    chip_w = Inches(1.65); chip_h = Inches(0.65)
    gap = Inches(0.10)
    total_w = chip_w * len(phases) + gap * (len(phases) - 1)
    start_x = (SLIDE_W - total_w) // 2
    y = Inches(2.5)
    for i, (label, col) in enumerate(phases):
        x = start_x + (chip_w + gap) * i
        _phase_chip(s, x, y, chip_w, chip_h, label, col)
        if i < len(phases) - 1:
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        x + chip_w - Inches(0.02),
                                        y + Inches(0.20),
                                        gap + Inches(0.04), Inches(0.25))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = TEXT_LIGHT
            arrow.line.fill.background()
    # Phase highlights
    _add_text(s, Inches(0.5), Inches(3.6), Inches(12), Inches(0.4),
              "Each phase shows a custom hero banner, an agent grid, an activity stream, and human approval gates.",
              size=14, color=TEXT_MID, align=PP_ALIGN.CENTER)
    # Bottom row of details
    details = [
        ("BRD extraction", "AI parses & classifies", BRAND_BLUE),
        ("Architecture", "C4 + API + ERD", BRAND_INDIGO),
        ("Real code gen", "writes bank app", RGBColor(0x08, 0x91, 0xB2)),
        ("Multi-port test", "load + UAT", BRAND_GREEN),
        ("Hardening", "SAST + secrets", BRAND_RED),
        ("Auto deploy", ":3001 & :3002", BRAND_PURPLE),
        ("Close-out", "post-mortem", BRAND_AMBER),
    ]
    y2 = Inches(4.4)
    for i, (title, sub, col) in enumerate(details):
        x = start_x + (chip_w + gap) * i
        _add_text(s, x, y2, chip_w, Inches(0.3), title, size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
        _add_text(s, x, y2 + Inches(0.32), chip_w, Inches(0.4), sub,
                  size=10, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    _add_footer(s, idx, total)


# ── Slide 4: Requirements features ────────────────────────────────
def slide_requirements(idx, total):
    s = prs.slides.add_slide(BLANK)
    _page_bg(s)
    _header(s, "PHASE · REQUIREMENTS", "Real AI requirement extraction",
            "Paste your project description — the requirement agent parses it and produces a structured BRD.")
    _card(s, Inches(0.5), Inches(2.3), Inches(6.0), Inches(4.4),
          "Heuristic NLP parser",
          ["Modal verbs → priority — must / should / nice",
           "Port detection — `port 3001`, `PORT=3001`, `:3001`",
           "Risk count — security keywords (auth, jwt, encrypt, …)",
           "Stakeholder patterns — PO, Eng Lead, QA, Architect",
           "Action verbs — login, register, validate, view, generate",
           "FR-001 … FR-NNN auto-numbering",
           "Works offline — no LLM dependency, no API key"], accent=BRAND_BLUE)
    _card(s, Inches(6.8), Inches(2.3), Inches(6.0), Inches(4.4),
          "Two-step input + extract",
          ["Paste tab — textarea with live char-count",
           "Upload tab — TXT / MD via FileReader.readAsText",
           "Validate Requirement button kicks off /api/extract",
           "Animated 4-step progress while the API call runs",
           "Extract card — stats, port grid, requirement list",
           "Stakeholders shown as compact pills",
           "Edit input button to re-paste & re-extract"], accent=BRAND_INDIGO)
    _add_footer(s, idx, total)


# ── Slide 5: Project type + tech stack ────────────────────────────
def slide_project_choice(idx, total):
    s = prs.slides.add_slide(BLANK)
    _page_bg(s)
    _header(s, "PHASE · REQUIREMENTS", "Project type & technology stack",
            "After extraction, choose how to deliver — branch from existing repo or scaffold a fresh stack.")
    _card(s, Inches(0.5), Inches(2.3), Inches(6.0), Inches(2.0),
          "New project (greenfield)",
          ["Three-option tech stack picker:",
           "Java / J2EE — Spring · JPA · Tomcat",
           "Python — FastAPI / Django · pip",
           ".NET — C# · ASP.NET · NuGet"], accent=BRAND_BLUE)
    _card(s, Inches(0.5), Inches(4.45), Inches(6.0), Inches(2.0),
          "Existing project (branch)",
          ["Auto-creates feature/abc-bank-core branch",
           "PR agent panel appears below",
           "Editable branch name input",
           "Triggers a real Jenkins pipeline run"], accent=BRAND_INDIGO)
    _card(s, Inches(6.8), Inches(2.3), Inches(6.0), Inches(4.4),
          "Gating logic",
          ["Run phases stays disabled until choices are made",
           "Helper hints — \"↑ Select project type to continue\"",
           "Switching project type auto-clears tech stack",
           "Choices persist across the whole pipeline",
           "Activity log records every selection",
           "All choices reset on Hard Reset"], accent=BRAND_PURPLE)
    _add_footer(s, idx, total)


# ── Slide 6: Development — real bank app ──────────────────────────
def slide_dev(idx, total):
    s = prs.slides.add_slide(BLANK)
    _page_bg(s)
    _header(s, "PHASE · DEVELOPMENT", "Real bank app, real running ports",
            "When dev-phase agents finish, the orchestrator writes a working ABC Bank app to disk and launches it.")
    _card(s, Inches(0.5), Inches(2.3), Inches(6.0), Inches(4.4),
          "Generated files (sdlc-orchestrator/generated-app/)",
          ["index.html — login page (demo / demo123)",
           "dashboard.html — balance + transactions",
           "data.js — users + sample transactions",
           "app.js — auth, render, logout (sessionStorage)",
           "styles.css — light theme matching orchestrator",
           "README.md — run + demo credentials",
           "All files written by /api/generate-app endpoint"], accent=BRAND_INDIGO)
    _card(s, Inches(6.8), Inches(2.3), Inches(6.0), Inches(4.4),
          "Multi-port launcher",
          ["Launcher panel appears after dev agents complete",
           ":3001 and :3002 buttons — click Launch",
           "POST /api/launch?port=N spawns python -m http.server",
           "Real subprocess.Popen with PID tracking",
           "Click links → opens the real bank app in browser",
           "Stop button kills the subprocess cleanly",
           "Footer in app shows actual port number"], accent=BRAND_PURPLE)
    _add_footer(s, idx, total)


# ── Slide 7: Jenkins integration ──────────────────────────────────
def slide_jenkins(idx, total):
    s = prs.slides.add_slide(BLANK)
    _page_bg(s)
    _header(s, "INTEGRATION", "Jenkins pipeline trigger",
            "PR agent triggers a real Jenkins build (or a built-in mock when credentials aren't set).")
    _card(s, Inches(0.5), Inches(2.3), Inches(6.0), Inches(4.4),
          "Real Jenkins mode",
          ["JENKINS_URL / USER / TOKEN env vars",
           "POST /job/<job>/buildWithParameters?BRANCH=…",
           "HTTP Basic auth headers sent server-side",
           "Polls every 2.5s — queue → build number → stages",
           "Pipeline Stage View plugin → live stage colors",
           "\"View in Jenkins\" link to the build page",
           "Detected real Jenkins 2.529 on :8080 in testing"], accent=BRAND_GREEN)
    _card(s, Inches(6.8), Inches(2.3), Inches(6.0), Inches(4.4),
          "Auto-mock fallback",
          ["No creds set → mock kicks in automatically",
           "5 stages: Checkout SCM → Build → Test → Deploy → Notify",
           "~2.5s per stage · brief queued phase",
           "Same UI / same polling / same code path",
           "Amber MOCK badge on the build footer",
           "\"View in Jenkins\" link hidden for mock URLs",
           "Lets the full demo run without any setup"], accent=BRAND_AMBER)
    _add_footer(s, idx, total)


# ── Slide 8: Sign-off + reject feature ────────────────────────────
def slide_signoff(idx, total):
    s = prs.slides.add_slide(BLANK)
    _page_bg(s)
    _header(s, "GOVERNANCE", "Sign-off, reject reason & verify",
            "Human gates between every phase — with auditable rejection trail.")
    _card(s, Inches(0.5), Inches(2.3), Inches(6.0), Inches(4.4),
          "Stakeholder sign-off",
          ["Multiple approvers per phase (PO, Eng, QA, Security)",
           "Each reviewer signs independently",
           "Comments shown inline after approval",
           "Approve & advance unlocks only when all signed",
           "Auto-advances to the next phase on approval",
           "Activity log records every signature"], accent=BRAND_BLUE)
    _card(s, Inches(6.8), Inches(2.3), Inches(6.0), Inches(4.4),
          "Reject → reason → verify → proceed",
          ["Click Reject → inline textarea for reason",
           "Confirm rejection stores reason in state",
           "Rejected banner shows the reason text",
           "☐ \"Reason addressed\" verification checkbox",
           "Run phases / Back to <prev> disabled until checked",
           "Approval auto-clears the rejection record"], accent=BRAND_RED)
    _add_footer(s, idx, total)


# ── Slide 9: Activity log + observability ─────────────────────────
def slide_observability(idx, total):
    s = prs.slides.add_slide(BLANK)
    _page_bg(s)
    _header(s, "OBSERVABILITY", "Real-time activity log",
            "Every action — agent start, approval, rejection, build status — appears in the live log.")
    _card(s, Inches(0.5), Inches(2.3), Inches(6.0), Inches(4.4),
          "What gets logged",
          ["Phase transitions (active → running → pending → done)",
           "Each AI agent's stream lines as they execute",
           "Stakeholder signatures",
           "Rejection reasons + verification toggles",
           "BRD extraction counts (requirements, ports, risks)",
           "Bank app generation + launch / stop",
           "Jenkins build queue + stage + result"], accent=BRAND_BLUE)
    _card(s, Inches(6.8), Inches(2.3), Inches(6.0), Inches(4.4),
          "Log entry styling",
          ["Timestamps in monospace (JetBrains Mono)",
           "Severity tags: INFO · SUCCESS · WARN · DANGER",
           "Color-coded pills per severity",
           "Newest entries on top",
           "200-entry rolling window",
           "Survives phase navigation"], accent=BRAND_AMBER)
    _add_footer(s, idx, total)


# ── Slide 10: Phase visual identity ───────────────────────────────
def slide_visual(idx, total):
    s = prs.slides.add_slide(BLANK)
    _page_bg(s)
    _header(s, "DESIGN", "Phase-specific visual identity",
            "Each phase opens with its own gradient hero banner — instant context for where you are in the pipeline.")
    bands = [
        ("REQUIREMENTS", BRAND_BLUE, BRAND_INDIGO),
        ("DESIGN", BRAND_INDIGO, BRAND_PURPLE),
        ("DEVELOPMENT", RGBColor(0x08, 0x91, 0xB2), BRAND_BLUE),
        ("TESTING", BRAND_GREEN, RGBColor(0x08, 0x91, 0xB2)),
        ("SECURITY", BRAND_RED, RGBColor(0xEA, 0x58, 0x0C)),
        ("DEPLOYMENT", BRAND_PURPLE, BRAND_BLUE),
        ("REVIEW", BRAND_AMBER, BRAND_RED),
    ]
    y = Inches(2.4)
    band_h = Inches(0.55)
    band_w = Inches(12.3)
    for i, (label, c1, c2) in enumerate(bands):
        _add_rounded(s, Inches(0.5), y + band_h * i + Inches(0.05) * i, band_w, band_h, c1, radius=0.3)
        _add_rounded(s, Inches(0.5) + band_w / 2, y + band_h * i + Inches(0.05) * i, band_w / 2, band_h, c2, radius=0.3)
        _add_text(s, Inches(1.0), y + band_h * i + Inches(0.05) * i + Inches(0.13),
                  Inches(11), Inches(0.4), label, size=14, bold=True, color=BG_WHITE)
    _add_footer(s, idx, total)


# ── Slide 11: Tech stack used ─────────────────────────────────────
def slide_techstack(idx, total):
    s = prs.slides.add_slide(BLANK)
    _page_bg(s)
    _header(s, "BUILT WITH", "Technology stack",
            "Two interoperable frontends + one minimal Python backend + Jenkins integration.")
    _card(s, Inches(0.5), Inches(2.3), Inches(4.0), Inches(4.4),
          "Frontend — vanilla",
          ["HTML5 + CSS3 (10 modular stylesheets)",
           "Vanilla JS — ES6+ IIFE modules",
           "Tabler Icons webfont",
           "Inter + JetBrains Mono fonts",
           "No build step — pure script tags",
           "Served by Python http.server at :3000"], accent=BRAND_BLUE)
    _card(s, Inches(4.7), Inches(2.3), Inches(4.0), Inches(4.4),
          "Frontend — Next.js",
          ["Next.js 15 (App Router)",
           "React 19 client components",
           "TypeScript 5.7 (strict mode)",
           "Zustand 5 — single store",
           "Custom useAgentRunner hook",
           "Hot-reloads at :3030 · proxies /api/* to :3000"], accent=BRAND_INDIGO)
    _card(s, Inches(8.9), Inches(2.3), Inches(4.0), Inches(4.4),
          "Backend & integrations",
          ["Python 3.13 — stdlib only, no pip deps",
           "http.server + ThreadingTCPServer",
           "subprocess.Popen for bank app launches",
           "urllib + Basic auth → Jenkins REST",
           "Heuristic NLP via re (regex)",
           "Real Jenkins 2.529 + auto-mock fallback"], accent=BRAND_PURPLE)
    _add_footer(s, idx, total)


# ── Slide 12: Architecture diagram ────────────────────────────────
def slide_architecture(idx, total):
    s = prs.slides.add_slide(BLANK)
    _page_bg(s)
    _header(s, "ARCHITECTURE", "How the pieces fit together",
            "Browser → Next.js (proxies /api/*) → Python backend → real subprocesses + Jenkins.")
    # Browser box
    _add_rounded(s, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.4), BG_WHITE, line=BORDER, radius=0.05)
    _add_text(s, Inches(0.8), Inches(2.45), Inches(2), Inches(0.4), "BROWSER",
              size=11, bold=True, color=TEXT_LIGHT)
    boxes = [
        (":3000", "Vanilla orchestrator", BRAND_BLUE),
        (":3030", "Next.js orchestrator", BRAND_INDIGO),
        (":3001", "ABC Bank instance #1", BRAND_PURPLE),
        (":3002", "ABC Bank instance #2", BRAND_PURPLE),
        (":8080", "Jenkins (or mock)", BRAND_GREEN),
    ]
    for i, (port, label, col) in enumerate(boxes):
        x = Inches(0.8 + i * 2.4)
        _add_rounded(s, x, Inches(2.85), Inches(2.2), Inches(0.65), col, radius=0.15)
        _add_text(s, x, Inches(2.92), Inches(2.2), Inches(0.3), port,
                  size=12, bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)
        _add_text(s, x, Inches(3.18), Inches(2.2), Inches(0.3), label,
                  size=9, color=BG_WHITE, align=PP_ALIGN.CENTER)

    # arrows down
    for i in range(2):
        arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                    Inches(1.5 + i * 2.4), Inches(3.9),
                                    Inches(0.4), Inches(0.4))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = TEXT_LIGHT
        arrow.line.fill.background()

    # Python backend
    _add_rounded(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(1.4), BG_WHITE, line=BORDER, radius=0.05)
    _add_text(s, Inches(0.8), Inches(4.65), Inches(8), Inches(0.4),
              "PYTHON BACKEND  ·  server.py", size=11, bold=True, color=TEXT_LIGHT)
    apis = ["/api/extract", "/api/generate-app", "/api/launch + /stop", "/api/jenkins/*"]
    for i, label in enumerate(apis):
        x = Inches(0.8 + i * 3)
        _add_rounded(s, x, Inches(5.0), Inches(2.8), Inches(0.7), RGBColor(0xEE, 0xF2, 0xF7),
                     line=BORDER, radius=0.15)
        _add_text(s, x, Inches(5.18), Inches(2.8), Inches(0.4), label,
                  size=12, bold=True, color=BRAND_BLUE, align=PP_ALIGN.CENTER, font="Consolas")

    # bottom note
    _add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4),
              "Same backend serves both frontends · subprocess.Popen owns the bank-app instances · Jenkins is optional (mock auto-fallback)",
              size=12, color=TEXT_MID, align=PP_ALIGN.CENTER)
    _add_footer(s, idx, total)


# ── Slide 13: How to run ──────────────────────────────────────────
def slide_run(idx, total):
    s = prs.slides.add_slide(BLANK)
    _page_bg(s)
    _header(s, "RUN IT", "Two PowerShell commands",
            "Python backend + Next.js dev server. Optional Jenkins env vars otherwise it mocks.")
    # Terminal 1
    _add_text(s, Inches(0.5), Inches(2.3), Inches(6), Inches(0.4),
              "Terminal 1 — Python backend (:3000)", size=14, bold=True, color=TEXT_DARK)
    _add_rounded(s, Inches(0.5), Inches(2.75), Inches(6), Inches(1.6), TEXT_DARK, radius=0.05)
    _add_text(s, Inches(0.7), Inches(2.9), Inches(5.6), Inches(1.4),
              "cd sdlc-orchestrator\npython server.py",
              size=14, color=BG_WHITE, font="Consolas")

    # Terminal 2
    _add_text(s, Inches(6.8), Inches(2.3), Inches(6), Inches(0.4),
              "Terminal 2 — Next.js dev server (:3030)", size=14, bold=True, color=TEXT_DARK)
    _add_rounded(s, Inches(6.8), Inches(2.75), Inches(6), Inches(1.6), TEXT_DARK, radius=0.05)
    _add_text(s, Inches(7.0), Inches(2.9), Inches(5.6), Inches(1.4),
              "cd sdlc-orchestrator-next\nnpm run dev",
              size=14, color=BG_WHITE, font="Consolas")

    # Optional Jenkins
    _add_text(s, Inches(0.5), Inches(4.6), Inches(12), Inches(0.4),
              "Optional — connect to real Jenkins (otherwise mock runs automatically)",
              size=14, bold=True, color=TEXT_DARK)
    _add_rounded(s, Inches(0.5), Inches(5.05), Inches(12.3), Inches(1.4), TEXT_DARK, radius=0.05)
    _add_text(s, Inches(0.7), Inches(5.2), Inches(12), Inches(1.2),
              "$env:JENKINS_USER  = \"<your-username>\"\n"
              "$env:JENKINS_TOKEN = \"<your-api-token>\"\n"
              "$env:JENKINS_JOB   = \"abc-bank\"",
              size=14, color=BG_WHITE, font="Consolas")
    _add_footer(s, idx, total)


# ── Slide 14: Closing ─────────────────────────────────────────────
def slide_thanks(idx, total):
    s = prs.slides.add_slide(BLANK)
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BRAND_BLUE)
    _add_rect(s, Inches(7), 0, SLIDE_W - Inches(7), SLIDE_H, BRAND_PURPLE)
    _add_text(s, Inches(0.5), Inches(2.5), Inches(12), Inches(1.5),
              "Thank you", size=72, bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(0.5), Inches(4.0), Inches(12), Inches(0.6),
              "AI Orchestrator · AI-Powered SDLC Pipeline",
              size=22, color=BG_WHITE, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(0.5), Inches(5.0), Inches(12), Inches(0.5),
              "Questions?", size=18, color=BG_WHITE, align=PP_ALIGN.CENTER)


# ── Build deck ────────────────────────────────────────────────────
builders = [
    slide_title,
    slide_what,
    slide_pipeline,
    slide_requirements,
    slide_project_choice,
    slide_dev,
    slide_jenkins,
    slide_signoff,
    slide_observability,
    slide_visual,
    slide_techstack,
    slide_architecture,
    slide_run,
    slide_thanks,
]
total = len(builders)
for i, fn in enumerate(builders, start=1):
    fn(i, total)

OUT = "AI-Orchestrator-Features.pptx"
prs.save(OUT)
print(f"✓ wrote {OUT}  ({total} slides)")
